"""
Socratic-AI Presentation - Clean version with diagram placeholders
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# Team Info
PRESENTATION_DATA = {
    "title": "Socratic - AI Tutorial Agent",
    "subtitle": "An Intelligent Tutoring System using LangGraph & RAG",
    "guide": "Project Guide Name",
    "members": [
        ("Shreyas Panda", "PRN: XXXXXXXXXXXX"),
        ("Shreyas Vichare", "PRN: XXXXXXXXXXXX"),
        ("Pratik Nagargoje", "PRN: XXXXXXXXXXXX"),
    ],
}

# Colors
COLORS = {
    "primary": RGBColor(0x00, 0x52, 0x99),
    "secondary": RGBColor(0x00, 0x7A, 0xCC),
    "accent": RGBColor(0x00, 0x96, 0x88),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "text_dark": RGBColor(0x2D, 0x3A, 0x4A),
    "bg_light": RGBColor(0xF8, 0xF9, 0xFA),
    "placeholder": RGBColor(0xE0, 0xE0, 0xE0),
}

FONT = "Alegreya"


def set_font(run, size, bold=False, color=None):
    run.font.name = FONT
    run.font.size = size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_gradient_bg(slide):
    bg = slide.background.fill
    bg.gradient()
    bg.gradient_angle = 135
    bg.gradient_stops[0].color.rgb = RGBColor(0x0D, 0x47, 0xA1)
    bg.gradient_stops[1].color.rgb = RGBColor(0x19, 0x76, 0xD2)


def add_light_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["bg_light"]


def add_header(slide, prs, title):
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS["primary"]
    header.line.fill.background()
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), prs.slide_width - Inches(1), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_font(r, Pt(28), bold=True, color=COLORS["white"])


def add_bullets(slide, bullets, top=1.3, numbered=False):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        prefix = f"{i+1}. " if numbered else "•  "
        r.text = prefix + text
        set_font(r, Pt(16), color=COLORS["text_dark"])
        p.space_before = Pt(8)
        p.space_after = Pt(8)


def add_diagram_placeholder(slide, label, top=1.5, height=4.5):
    """Add a placeholder box for diagram"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(top), Inches(9), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS["placeholder"]
    box.line.color.rgb = RGBColor(0xBD, 0xBD, 0xBD)
    box.line.width = Pt(2)
    
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"\n\n\n[INSERT {label} HERE]"
    set_font(r, Pt(18), bold=True, color=RGBColor(0x75, 0x75, 0x75))


def add_table(slide, headers, rows, top=1.3):
    cols = len(headers)
    num_rows = len(rows) + 1
    
    table = slide.shapes.add_table(num_rows, cols, Inches(0.5), Inches(top), Inches(9), Inches(num_rows * 0.5)).table
    
    # Header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["primary"]
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["white"]
        p.font.name = FONT
    
    # Data
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["white"] if ri % 2 == 0 else RGBColor(0xE3, 0xF2, 0xFD)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS["text_dark"]
            p.font.name = FONT


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # SLIDE 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)
    
    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = PRESENTATION_DATA["title"]
    set_font(r, Pt(44), bold=True, color=COLORS["white"])
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(9), Inches(0.5))
    p2 = tb2.text_frame.paragraphs[0]
    r2 = p2.add_run()
    r2.text = PRESENTATION_DATA["subtitle"]
    set_font(r2, Pt(22), color=RGBColor(0xE0, 0xE0, 0xE0))
    p2.alignment = PP_ALIGN.CENTER
    
    # Guide
    tb3 = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(0.8))
    tf3 = tb3.text_frame
    p3a = tf3.paragraphs[0]
    r3a = p3a.add_run()
    r3a.text = "Under the Guidance of"
    set_font(r3a, Pt(14), color=RGBColor(0xE0, 0xE0, 0xE0))
    p3a.alignment = PP_ALIGN.CENTER
    
    p3b = tf3.add_paragraph()
    r3b = p3b.add_run()
    r3b.text = PRESENTATION_DATA["guide"]
    set_font(r3b, Pt(20), bold=True, color=COLORS["white"])
    p3b.alignment = PP_ALIGN.CENTER
    
    # Team
    tb4 = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(1.8))
    tf4 = tb4.text_frame
    p4a = tf4.paragraphs[0]
    r4a = p4a.add_run()
    r4a.text = "Submitted By"
    set_font(r4a, Pt(14), color=RGBColor(0xE0, 0xE0, 0xE0))
    p4a.alignment = PP_ALIGN.CENTER
    
    for name, prn in PRESENTATION_DATA["members"]:
        p4 = tf4.add_paragraph()
        r4 = p4.add_run()
        r4.text = f"{name}  |  {prn}"
        set_font(r4, Pt(16), color=COLORS["white"])
        p4.alignment = PP_ALIGN.CENTER
    
    print("  [1/17] Title")
    
    # SLIDE 2: Abstract
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_header(slide, prs, "Abstract")
    add_bullets(slide, [
        "An interactive intelligent tutoring system that teaches any subject using the Socratic Method.",
        "Built with Flask (Python), LangGraph for AI orchestration, and modern HTML/CSS/JavaScript frontend.",
        "Implements guided questioning to help learners discover answers themselves.",
        "Supports multi-modal input: text, voice (Speech-to-Text), and image analysis.",
        "Features RAG (Retrieval Augmented Generation) for personalized learning from user documents.",
        "Provides real-time streaming responses using Groq API with Llama 3.3 70B model."
    ])
    print("  [2/17] Abstract")
    
    # SLIDE 3: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_header(slide, prs, "Problem Statement")
    add_bullets(slide, [
        "Traditional education relies on passive learning methods without active student engagement.",
        "Students lack access to personalized tutoring that adapts to their individual learning pace.",
        "Existing AI chatbots provide direct answers instead of fostering critical thinking skills.",
        "Affordable 24/7 tutoring solutions are not available for all subjects and learners.",
        "Students need a system that can learn from their own study materials."
    ])
    print("  [3/17] Problem Statement")
    
    # SLIDE 4: Proposed Solution
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_header(slide, prs, "Proposed Solution")
    add_bullets(slide, [
        "Socratic-AI: An intelligent tutor using the Socratic method for guided learning.",
        "Generates adaptive, structured tutorials on any topic entered by the user.",
        "Supports multi-modal interaction: text input, voice recording, image upload.",
        "RAG-powered personal knowledge base from user-uploaded PDF and TXT documents.",
        "Real-time streaming responses using Groq API for fast inference.",
        "Progress tracking dashboard with learning statistics and session history."
    ])
    print("  [4/17] Proposed Solution")
    
    # SLIDE 5: Working Methodology (Diagram placeholder)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_header(slide, prs, "Working Methodology")
    add_diagram_placeholder(slide, "WORKING METHODOLOGY FLOWCHART")
    print("  [5/17] Working Methodology (placeholder)")
    
    # SLIDE 6: System Architecture (Diagram placeholder)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_header(slide, prs, "System Architecture")
    add_diagram_placeholder(slide, "SYSTEM ARCHITECTURE BLOCK DIAGRAM")
    print("  [6/17] System Architecture (placeholder)")
    
    # SLIDE 7: Libraries Used
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_header(slide, prs, "Libraries Used")
    add_table(slide, 
        ["Library", "Purpose"],
        [
            ["Flask", "Python web framework for backend API"],
            ["LangGraph", "State machine for AI agent orchestration"],
            ["LangChain", "LLM integration and chain building"],
            ["FAISS", "Vector similarity search for RAG"],
            ["HuggingFace", "Embedding model (all-MiniLM-L6-v2)"],
            ["Groq API", "Fast LLM inference (Llama 3.3 70B)"],
            ["Web Speech API", "Browser-based speech recognition"],
            ["Pillow", "Image processing and format conversion"],
            ["SQLite", "Session and progress data storage"],
        ]
    )
    print("  [7/17] Libraries Used")
    
    # SLIDE 8: RAG Implementation (Diagram placeholder)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_header(slide, prs, "RAG Implementation")
    add_diagram_placeholder(slide, "RAG FLOW DIAGRAM")
    print("  [8/17] RAG Implementation (placeholder)")
    
    # SLIDE 9: User Interaction Flow (Diagram placeholder)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_header(slide, prs, "User Interaction Flow")
    add_diagram_placeholder(slide, "USER INTERACTION FLOWCHART")
    print("  [9/17] User Interaction Flow (placeholder)")
    
    # SLIDE 10: Implementation Details
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_header(slide, prs, "Implementation Details")
    add_table(slide,
        ["File", "Purpose"],
        [
            ["app.py", "Flask routes, session management, API endpoints"],
            ["tutorial_agent.py", "LangGraph state machine, workflow nodes"],
            ["LLM_api.py", "Groq/OpenRouter API client configuration"],
            ["rag_engine.py", "RAG facade, document processing"],
            ["image_handler.py", "Image upload, vision model integration"],
            ["database.py", "SQLite operations, progress tracking"],
            ["main.js", "Frontend streaming, TTS, voice recording"],
        ]
    )
    print("  [10/17] Implementation Details")
    
    # SLIDE 11: Results
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_header(slide, prs, "Results")
    add_bullets(slide, [
        "Real-time streaming responses with stop generation capability implemented.",
        "Multi-modal input working seamlessly (text, voice, image).",
        "RAG-powered context-aware responses from uploaded documents verified.",
        "Modern dark theme UI (Midnight Pro) for comfortable extended usage.",
        "Progress dashboard with learning statistics and session history.",
        "PDF export functionality generating formatted lesson documents."
    ])
    print("  [11/17] Results")
    
    # SLIDE 12: Advantages
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_header(slide, prs, "Advantages")
    add_bullets(slide, [
        "Active Learning: Socratic method fosters critical thinking and deeper understanding.",
        "High Performance: Groq API provides near-instant response generation.",
        "Personalized: RAG allows learning from user's own study materials.",
        "Accessible: Web-based application works on any device with modern browser.",
        "Cost-Effective: Uses free-tier APIs (Groq: 14,400 requests/day).",
        "Privacy-Focused: Local processing with no external data storage."
    ])
    print("  [12/17] Advantages")
    
    # SLIDE 13: Limitations
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_header(slide, prs, "Limitations")
    add_bullets(slide, [
        "Requires stable internet connection for LLM API calls.",
        "Voice recognition accuracy depends on browser support and audio quality.",
        "Large document uploads may slow down RAG processing time.",
        "API rate limits on free tier (Groq: 14,400/day).",
        "Image analysis limited to supported formats (JPEG, PNG, HEIC).",
        "Currently single-user focused without multi-user authentication."
    ])
    print("  [13/17] Limitations")
    
    # SLIDE 14: Applications
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_header(slide, prs, "Applications")
    add_bullets(slide, [
        "Self-paced learning for students of all levels.",
        "Exam preparation using personal study materials.",
        "Concept clarification through Socratic dialogue.",
        "Document-based Q&A for research and reference.",
        "Accessible tutoring for remote and rural learners.",
        "Supplementary learning tool for classroom education."
    ])
    print("  [14/17] Applications")
    
    # SLIDE 15: Future Scope
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_header(slide, prs, "Future Scope")
    add_bullets(slide, [
        "User Authentication: Multi-user support with persistent learning history.",
        "Mobile Applications: Native iOS and Android apps.",
        "Video Analysis: Support for educational video content.",
        "Multi-Language Support: Regional language interface and responses.",
        "Advanced Analytics: Learning pattern analysis and recommendations.",
        "Collaborative Learning: Study groups and shared sessions.",
        "LMS Integration: Connect with Moodle, Canvas, Blackboard."
    ])
    print("  [15/17] Future Scope")
    
    # SLIDE 16: References
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_light_bg(slide)
    add_header(slide, prs, "References")
    add_bullets(slide, [
        "LangGraph Documentation - https://langchain-ai.github.io/langgraph/",
        "Groq API Documentation - https://console.groq.com/docs/",
        "FAISS Library - https://github.com/facebookresearch/faiss",
        "Flask Documentation - https://flask.palletsprojects.com/",
        "Web Speech API - MDN Web Docs",
        "HuggingFace Transformers - https://huggingface.co/docs/transformers/",
        "OpenRouter API - https://openrouter.ai/docs"
    ], numbered=True)
    print("  [16/17] References")
    
    # SLIDE 17: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)
    
    tb = slide.shapes.add_textbox(Inches(0), Inches(2.5), Inches(10), Inches(1))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Thank You"
    set_font(r, Pt(52), bold=True, color=COLORS["white"])
    p.alignment = PP_ALIGN.CENTER
    
    tb2 = slide.shapes.add_textbox(Inches(0), Inches(4.0), Inches(10), Inches(0.5))
    p2 = tb2.text_frame.paragraphs[0]
    names = [m[0] for m in PRESENTATION_DATA["members"]]
    r2 = p2.add_run()
    r2.text = "  |  ".join(names)
    set_font(r2, Pt(18), color=RGBColor(0xE0, 0xE0, 0xE0))
    p2.alignment = PP_ALIGN.CENTER
    
    tb3 = slide.shapes.add_textbox(Inches(0), Inches(5.0), Inches(10), Inches(0.5))
    p3 = tb3.text_frame.paragraphs[0]
    r3 = p3.add_run()
    r3.text = "Questions?"
    set_font(r3, Pt(24), color=COLORS["accent"])
    p3.alignment = PP_ALIGN.CENTER
    
    print("  [17/17] Thank You")
    
    # Save
    output = "Socratic_AI_Presentation.pptx"
    prs.save(output)
    print(f"\nSaved: {output}")
    print("\nDiagram placeholders on slides: 5, 6, 8, 9")
    print("Replace the gray boxes with your Mermaid diagrams!")
    return output


if __name__ == "__main__":
    print("Creating presentation with diagram placeholders...\n")
    create_presentation()
